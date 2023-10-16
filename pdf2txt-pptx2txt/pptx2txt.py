import os
from pptx import Presentation

# Replace this with the directory path where your PowerPoint files are located
pptx_directory = '/Users/Jyanqa/Desktop/thesis/prerequisite-text-extracttion/pdf2txt-pptx2txt/pptx/'

# Replace 'output.txt' with the name of the TXT file you want to create
txt_file_path = '/Users/Jyanqa/Desktop/thesis/prerequisite-text-extracttion/pdf2txt-pptx2txt/'+ 'output_pptx.txt'

# Initialize an empty string to store the extracted text
text = ""

# Loop through files in the directory
for filename in os.listdir(pptx_directory):
    if filename.endswith('.pptx'):
        full_path = os.path.join(pptx_directory, filename)

        # Load the PowerPoint presentation
        presentation = Presentation(full_path)

        # Iterate through slides and extract text
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "

# Write the extracted text to the output file
with open(txt_file_path, 'a') as f:
    f.write(text)

# Close the output file
f.close()
